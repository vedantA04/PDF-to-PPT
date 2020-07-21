from pptx import Presentation
from pptx.util import Inches
import os
import natsort
from PIL import Image

def imgs2powerpoint(outpath):
    prs = Presentation()
    img = Image.open("Images/Picture1.png")
    width , height = img.size
    prs.slide_width = Inches(width / 500)
    prs.slide_height = Inches(height / 500)
    blank_slide_layout = prs.slide_layouts[6]
    imagelist = os.listdir("Images/")
    for image in natsort.natsorted(imagelist):
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = 0
        slide.shapes.add_picture("Images/{}".format(image), left, top, height = prs.slide_height, width = prs.slide_width)
    prs.save(outpath)
